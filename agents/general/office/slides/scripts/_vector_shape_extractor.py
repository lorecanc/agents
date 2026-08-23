"""
OOXML vector shape properties extraction (fill, stroke, rotation, geometry).

Reads shape XML (``p:sp``, ``p:cxnSp``) and produces structured dicts for
``immutable_shapes`` entries in the design system ``index.json``.

*Fill* — ``a:solidFill`` (rgb / schemeClr), ``a:noFill``.
*Stroke* — ``a:ln`` width, colour, dash pattern, arrows.
*Rotation* — ``a:xfrm/@rot`` (centiseconds → degrees).
*Geometry* — ``a:prstGeom`` (preset name) or ``a:custGeom`` (SVG path string).
*Group shapes* — recursively flattens children with offset-adjusted coords.
"""

from __future__ import annotations

from typing import Any

from lxml import etree

from _ooxml_utils import NS, emu_to_px, emu_to_pt, resolve_scheme


# ── Helper: get spPr from any shape element ────────────────────────────

def _get_spPr(element: etree._Element) -> etree._Element | None:
    """Return ``p:spPr`` or ``a:spPr`` from a shape XML element."""
    for xpath in ("p:spPr", "a:spPr", ".//p:spPr", ".//a:spPr"):
        result = element.find(xpath, NS)
        if result is not None:
            return result
    return None


# ── Fill extraction ─────────────────────────────────────────────────────

def extract_fill(
    spPr: etree._Element,
    theme_colors: dict[str, str] | None = None,
) -> dict[str, Any] | None:
    """Extract fill info from ``p:spPr/a:spPr``.

    Returns ``None`` when no fill element is found (theme-inherited).
    Returns a dict with ``type`` and resolved ``color`` whenever possible.
    """
    theme_colors = theme_colors or {}
    fill_el = spPr.find("a:solidFill", NS)
    if fill_el is not None:
        for child in fill_el:
            tag = etree.QName(child.tag).localname
            if tag == "srgbClr":
                val = child.get("val")
                if val:
                    return {"type": "solid", "color": f"#{val}"}
            elif tag == "schemeClr":
                scheme = child.get("val", "")
                hex_c = resolve_scheme(scheme, theme_colors) or ""
                result: dict[str, Any] = {"type": "solid", "scheme": scheme}
                if hex_c:
                    result["color"] = hex_c
                return result
            elif tag == "sysClr":
                val = child.get("lastClr") or child.get("val", "")
                return {"type": "solid", "color": f"#{val}"}
            break
        # solidFill but no recognized child
        return {"type": "solid"}

    nofill = spPr.find("a:noFill", NS)
    if nofill is not None:
        return {"type": "noFill"}

    grad = spPr.find("a:gradFill", NS)
    if grad is not None:
        return {"type": "gradient"}

    blip = spPr.find("a:blipFill", NS)
    if blip is not None:
        return {"type": "picture"}

    return None  # inherited / none found


# ── Stroke extraction ───────────────────────────────────────────────────

def extract_stroke(
    spPr: etree._Element,
    theme_colors: dict[str, str] | None = None,
) -> dict[str, Any] | None:
    """Extract stroke (line) properties from ``a:ln``."""
    theme_colors = theme_colors or {}
    ln = spPr.find("a:ln", NS)
    if ln is None:
        return None

    w_str = ln.get("w")
    width_pt = emu_to_pt(int(w_str)) if w_str else None

    # Stroke colour
    stroke_fill: dict[str, Any] = {}
    solid = ln.find("a:solidFill", NS)
    if solid is not None:
        for child in solid:
            tag = etree.QName(child.tag).localname
            if tag == "srgbClr":
                val = child.get("val")
                if val:
                    stroke_fill["color"] = f"#{val}"
            elif tag == "schemeClr":
                scheme = child.get("val", "")
                stroke_fill["scheme"] = scheme
                hex_c = resolve_scheme(scheme, theme_colors) or ""
                if hex_c:
                    stroke_fill["color"] = hex_c
            elif tag == "sysClr":
                val = child.get("lastClr") or child.get("val", "")
                stroke_fill["color"] = f"#{val}"
            break

    nofill = ln.find("a:noFill", NS)
    if nofill is not None:
        stroke_fill["type"] = "noFill"
    elif stroke_fill:
        stroke_fill["type"] = "solid"

    # Dash pattern
    dash = ln.find("a:prstDash", NS)
    dash_val = dash.get("val") if dash is not None else None

    # Arrow heads
    head = ln.find("a:headEnd", NS)
    tail = ln.find("a:tailEnd", NS)

    result: dict[str, Any] = {}
    if stroke_fill:
        result.update(stroke_fill)
    if width_pt is not None:
        result["width"] = width_pt
    if dash_val:
        result["dash"] = dash_val
    if head is not None:
        result["head_end"] = {
            k: head.get(k) for k in ("type", "w", "len") if head.get(k)
        }
    if tail is not None:
        result["tail_end"] = {
            k: tail.get(k) for k in ("type", "w", "len") if tail.get(k)
        }

    return result if result else None


# ── Rotation extraction ──────────────────────────────────────────────────

def extract_transform(spPr: etree._Element) -> dict[str, Any]:
    """Extract rotation (degrees) and flips from ``a:xfrm``."""
    xfrm = spPr.find("a:xfrm", NS)
    if xfrm is None:
        return {"rotation": 0, "flip_h": False, "flip_v": False}

    rot_cs = xfrm.get("rot")  # centiseconds (1/100 of a degree)
    rotation = 0.0
    if rot_cs:
        rotation = round(int(rot_cs) / 60000, 1)  # centiseconds → degrees

    flip_h = xfrm.get("flipH") == "1"
    flip_v = xfrm.get("flipV") == "1"

    return {"rotation": rotation, "flip_h": flip_h, "flip_v": flip_v}


# ── Geometry extraction ──────────────────────────────────────────────────

_PRESET_SVG_GENERATORS: dict[str, str] = {
    # w, h will be substituted at render time
    "rect":        "<rect x=\"0\" y=\"0\" width=\"{w}\" height=\"{h}\"/>",
    "ellipse":     "<ellipse cx=\"{cx}\" cy=\"{cy}\" rx=\"{rx}\" ry=\"{ry}\"/>",
    "rtTriangle":  "<polygon points=\"{w},{h} 0,{h} 0,0\"/>",
    "triangle":    "<polygon points=\"{cx},{h} 0,0 {w},0\"/>",
}
"""Built-in SVG snippets for OOXML preset geometries.

Each template is resolved with ``w``, ``h`` (shape pixels), plus derived
values (``cx``, ``cy``, ``rx``, ``ry``).

Line geometry is handled separately by ``extract_preset_geometry``.
"""


def extract_preset_geometry(prst: etree._Element, w_px: int, h_px: int) -> dict[str, Any]:
    """Return a geometry dict for a ``a:prstGeom`` shape.

    ``w_px`` and ``h_px`` are the shape's pixel dimensions (pre-rotation).
    """
    name = prst.get("prst", "")
    result: dict[str, Any] = {"preset": name}

    # Special case: "line" preset — determine orientation from w/h
    # SVG snippet uses position 0 for both ends in the zero-dimension axis;
    # the SVG builder will adjust the SVG container width/height for visibility.
    if name == "line":
        if w_px >= h_px:
            snippet = f"<line x1=\"0\" y1=\"0\" x2=\"{w_px}\" y2=\"0\"/>"
        else:
            snippet = f"<line x1=\"0\" y1=\"0\" x2=\"0\" y2=\"{h_px}\"/>"
        result["svg_snippet"] = snippet
        return result

    svg_template = _PRESET_SVG_GENERATORS.get(name)
    if svg_template is not None:
        cx = round(w_px / 2)
        cy = round(h_px / 2)
        rx = round(w_px / 2)
        ry = round(h_px / 2)
        result["svg_snippet"] = svg_template.format(
            w=w_px, h=h_px, cx=cx, cy=cy, rx=rx, ry=ry
        )

    return result


def extract_custom_geometry(cust: etree._Element) -> dict[str, Any]:
    """Return a geometry dict for a ``a:custGeom`` shape.

    Converts the OOXML path commands into an SVG ``d`` string.
    """
    path = cust.find("a:pathLst/a:path", NS)
    if path is None:
        return {}

    path_w = path.get("w")
    path_h = path.get("h")

    svg_commands: list[str] = []
    for child in path:
        tag = etree.QName(child.tag).localname
        if tag == "moveTo":
            pt = child.find("a:pt", NS)
            if pt is not None:
                svg_commands.append(f"M {pt.get('x')} {pt.get('y')}")
        elif tag == "lnTo":
            pt = child.find("a:pt", NS)
            if pt is not None:
                svg_commands.append(f"L {pt.get('x')} {pt.get('y')}")
        elif tag == "cubicBezTo":
            pts = child.findall("a:pt", NS)
            if len(pts) == 3:
                svg_commands.append(
                    f"C {pts[0].get('x')} {pts[0].get('y')}"
                    f" {pts[1].get('x')} {pts[1].get('y')}"
                    f" {pts[2].get('x')} {pts[2].get('y')}"
                )
        elif tag == "quadBezTo":
            pts = child.findall("a:pt", NS)
            if len(pts) == 2:
                svg_commands.append(
                    f"Q {pts[0].get('x')} {pts[0].get('y')}"
                    f" {pts[1].get('x')} {pts[1].get('y')}"
                )
        elif tag == "arcTo":
            # Convert OOXML arcTo to approximate SVG arc
            # wR, hR = radii (in path coords); stAng, swAng in 1/60000 degrees
            wR = child.get("wR", "0")
            hR = child.get("hR", "0")
            # For skeleton purposes, approximate with a straight segment
            svg_commands.append(f"// arcTo wR={wR} hR={hR}")
        elif tag == "close":
            svg_commands.append("Z")

    result: dict[str, Any] = {}
    if svg_commands:
        result["path"] = " ".join(svg_commands)
    if path_w is not None:
        result["path_w"] = path_w
    if path_h is not None:
        result["path_h"] = path_h

    return result


def extract_geometry(
    spPr: etree._Element, w_px: int, h_px: int
) -> dict[str, Any]:
    """Extract geometry from a shape's spPr element."""
    prst = spPr.find("a:prstGeom", NS)
    if prst is not None:
        return extract_preset_geometry(prst, w_px, h_px)

    cust = spPr.find("a:custGeom", NS)
    if cust is not None:
        return extract_custom_geometry(cust)

    return {}


# ── Top-level extraction for a single shape ──────────────────────────────

def extract_vector_shape_properties(
    shape,
    theme_colors: dict[str, str] | None = None,
) -> dict[str, Any]:
    """Extract fill, stroke, rotation, and geometry from a layout/master shape.

    Returns a dict with keys: ``fill``, ``stroke``, ``rotation``, ``flip_h``,
    ``flip_v``, ``geometry``.

    Missing/inapplicable keys are omitted.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    theme_colors = theme_colors or {}
    xml_el = shape._element
    spPr = _get_spPr(xml_el)
    if spPr is None:
        return {}

    w_px = emu_to_px(shape.width)
    h_px = emu_to_px(shape.height)
    shape_type = shape.shape_type

    result: dict[str, Any] = {}

    # Fill
    if shape_type in (
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.FREEFORM,
    ):
        fill = extract_fill(spPr, theme_colors)
        if fill:
            result["fill"] = fill

    # Stroke (line properties) — relevant for all shape types
    stroke = extract_stroke(spPr, theme_colors)
    if stroke:
        result["stroke"] = stroke

    # Transform
    xform = extract_transform(spPr)
    if xform.get("rotation", 0) != 0:
        result["rotation"] = xform["rotation"]
    if xform.get("flip_h"):
        result["flip_h"] = True
    if xform.get("flip_v"):
        result["flip_v"] = True

    # Geometry
    geom = extract_geometry(spPr, w_px, h_px)
    if geom:
        result["geometry"] = geom

    return result


# ── Group fill inheritance ──────────────────────────────────────────────

def _has_grp_fill(spPr: etree._Element | None) -> bool:
    """Check if a shape uses ``a:grpFill`` (inherits fill from parent group)."""
    if spPr is None:
        return False
    return spPr.find("a:grpFill", NS) is not None


def _extract_group_spPr(group_shape) -> etree._Element | None:
    """Return the ``a:grpSpPr`` or ``p:grpSpPr`` element from a group shape."""
    return group_shape._element.find("a:grpSpPr", NS) or group_shape._element.find("p:grpSpPr", NS)


def _extract_group_inherited_fill(
    group_shape,
    theme_colors: dict[str, str],
    outer_inherited_fill: dict[str, Any] | None = None,
) -> dict[str, Any] | None:
    """Resolve the effective fill inherited by children of a group.

    Walks the group's own ``grpSpPr`` — if it contains ``a:grpFill``,
    falls back to the outer inherited fill. Otherwise extracts the
    group's own fill.
    """
    grp_spPr = _extract_group_spPr(group_shape)
    if grp_spPr is None:
        return outer_inherited_fill

    if _has_grp_fill(grp_spPr):
        return outer_inherited_fill

    group_fill = extract_fill(grp_spPr, theme_colors)
    return group_fill or outer_inherited_fill


# ── Group shape recursion ────────────────────────────────────────────────

def flatten_group_shape(
    group_shape,
    theme_colors: dict[str, str] | None = None,
    parent_offset_x: int = 0,
    parent_offset_y: int = 0,
    inherited_fill: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    """Recursively extract child shapes from a GROUP shape.

    Each child's position is adjusted relative to the group's offset so that
    coordinates remain absolute within the slide layout.

    ``inherited_fill`` is the fill inherited from ancestor groups, used to
    resolve ``a:grpFill`` on child shapes.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    theme_colors = theme_colors or {}
    group_x = emu_to_px(group_shape.left)
    group_y = emu_to_px(group_shape.top)
    abs_x = parent_offset_x + group_x
    abs_y = parent_offset_y + group_y

    # Resolve the fill to pass to children
    child_inherited_fill = _extract_group_inherited_fill(
        group_shape, theme_colors, inherited_fill
    )

    children: list[dict[str, Any]] = []
    try:
        sub_shapes = list(group_shape.shapes)
    except Exception:
        return children

    for child in sub_shapes:
        child_type = child.shape_type
        w_px = emu_to_px(child.width)
        h_px = emu_to_px(child.height)
        x_px = abs_x + emu_to_px(child.left)
        y_px = abs_y + emu_to_px(child.top)

        if child_type == MSO_SHAPE_TYPE.GROUP:
            children.extend(
                flatten_group_shape(child, theme_colors, abs_x, abs_y, child_inherited_fill)
            )
        elif child_type in (
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.LINE,
        ):
            entry = _make_shape_entry(
                child, child_type, x_px, y_px, w_px, h_px, theme_colors,
                inherited_fill=child_inherited_fill,
            )
            if entry:
                children.append(entry)

    return children


def _make_shape_entry(
    shape,
    shape_type,
    x_px: int,
    y_px: int,
    w_px: int,
    h_px: int,
    theme_colors: dict[str, str],
    *,
    inherited_fill: dict[str, Any] | None = None,
) -> dict[str, Any] | None:
    """Build a single immutable_shapes entry for a non-PH shape.

    When ``inherited_fill`` is provided (from a parent group's ``grpSpPr``),
    it is used for shapes that use ``a:grpFill`` (group fill inheritance).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Label
    _SHAPE_TYPE_LABELS = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.LINE: "line",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
        MSO_SHAPE_TYPE.GROUP: "group",
    }
    _conn = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
    if _conn is not None:
        _SHAPE_TYPE_LABELS[_conn] = "connector"

    label = _SHAPE_TYPE_LABELS.get(shape_type, f"shape_{shape_type}")

    entry: dict[str, Any] = {
        "shape_type": label,
        "x": x_px,
        "y": y_px,
        "w": w_px,
        "h": h_px,
    }

    # Subtype (auto_shape_type) — only meaningful for auto_shapes
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            if hasattr(shape, "auto_shape_type"):
                subtype = str(shape.auto_shape_type).lower().split(".")[-1]
                if subtype and subtype != "mixed":
                    entry["subtype"] = subtype
        except Exception:
            pass

    # Vector properties (fill, stroke, rotation, geometry)
    # Skip for PICTURE shapes — they're extracted as images, not vectors
    if shape_type != MSO_SHAPE_TYPE.PICTURE:
        props = extract_vector_shape_properties(shape, theme_colors)

        # Handle a:grpFill — inherit fill from parent group chain
        if inherited_fill is not None:
            spPr = _get_spPr(shape._element)
            if _has_grp_fill(spPr):
                if inherited_fill:
                    props["fill"] = inherited_fill

        if props:
            entry.update(props)

    # Text
    if shape.has_text_frame:
        try:
            text = shape.text_frame.text.strip()
            if text:
                entry["text"] = text
        except Exception:
            pass

    return entry
