"""
Analyze a corporate .pptx template and produce a design system directory
with colors.md, typography.md, HTML layout skeletons, and a fingerprint
for idempotent caching.

Usage:
  python analyze_slide_master.py \\
    --input corporate.pptx \\
    --design-name acme-corp \\
    --output-dir /path/to/.opencode/office/slides/design/ \\
    [--force]
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import json
import os
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from lxml import etree  # type: ignore
from pptx import Presentation  # type: ignore

from _skeleton_builder import (
    classify_layout,
    generate_skeleton_html,
)
from _text_style_extractor import (
    extract_text_style,
    extract_master_text_styles,
)
from _vector_shape_extractor import (
    _make_shape_entry,
    flatten_group_shape,
)


# ── Theme font reference mapping ──────────────────────────────────────


def _resolve_theme_font(font_name: str, theme_fonts: dict[str, str]) -> str | None:
    """Resolve a ``+mj-lt`` / ``+mn-lt`` style theme font reference to a real
    font family name using the theme fonts dict (keys ``major``, ``minor``).
    """
    if not font_name or not font_name.startswith("+"):
        return font_name

    # "+mj-lt" → major-latin, "+mj-ea" → major-east-asian, etc.
    base = font_name[1:].lower()  # remove "+"
    if base.startswith("mj-"):
        return theme_fonts.get("major") or font_name
    if base.startswith("mn-"):
        return theme_fonts.get("minor") or font_name
    return font_name  # unknown theme reference — keep as-is


def _merge_master_text_style(
    text_style: dict[str, Any] | None,
    ph_type: str,
    master_text_styles: dict[str, Any],
    theme_fonts: dict[str, str],
) -> dict[str, Any] | None:
    """Backfill missing font and text properties in a per-shape ``text_style``
    from the slide master ``p:txStyles``.

    When a shape's own XML does not specify font_name, font_size_pt, or body_pr,
    those values are inherited from the corresponding master text style section:

    - TITLE / CTR_TITLE / SUBTITLE → master ``title`` section
    - BODY / VERTICAL_BODY → master ``body`` section
    - Everything else → master ``other`` section

    Theme font references (``+mj-lt``, ``+mn-lt``, etc.) are resolved to real
    font family names using ``theme_fonts``.
    """
    if not text_style:
        return text_style

    name_lower = (ph_type or "").lower()
    if name_lower in ("title", "ctr_title", "subtitle", "vertical_title"):
        section_key = "title"
    elif name_lower in ("body", "vertical_body"):
        section_key = "body"
    else:
        section_key = "other"

    master_section = master_text_styles.get(section_key, {})

    # Merge into level 1 run
    levels = text_style.setdefault("levels", {})
    master_level_1 = master_section.get("1", {})
    master_run = master_level_1.get("run", {})
    master_para = master_level_1.get("paragraph", {})

    lvl1 = levels.setdefault("1", {})
    if not isinstance(lvl1, dict):
        lvl1 = {}
        levels["1"] = lvl1

    # Run-level backfill
    run = lvl1.setdefault("run", {})
    if not isinstance(run, dict):
        run = {}
        lvl1["run"] = run

    for key in ("font_name", "font_name_ea", "font_name_cs", "font_size_pt",
                "bold", "italic", "kerning_pt"):
        if key not in run or run[key] is None:
            mv = master_run.get(key)
            if mv is not None:
                if key.startswith("font_name"):
                    mv = _resolve_theme_font(mv, theme_fonts)
                run[key] = mv

    # Also merge color if missing
    if "color" not in run or run["color"] is None:
        if "color" in master_run and master_run["color"] is not None:
            run["color"] = dict(master_run["color"])

    # Paragraph-level backfill
    para = lvl1.setdefault("paragraph", {})
    if not isinstance(para, dict):
        para = {}
        lvl1["paragraph"] = para

    for key in ("alignment", "line_spacing", "space_before", "space_after"):
        if key not in para or para[key] is None:
            mv = master_para.get(key)
            if mv is not None:
                para[key] = mv

    # Body-pr backfill
    master_body_pr: dict[str, Any] = {}
    # Master body styles don't have body_pr — those come from layout shapes.
    # We don't merge body_pr from master since master doesn't have it.

    return text_style


def _detect_layout_bg(layout, master, theme_colors: dict[str, str]) -> str | None:
    """Detect background colour of a slide layout.

    Tries in order:
    1. Layout background: srgbClr or schemeClr in solidFill
    2. Slide master background (same logic)
    3. Theme ``lt1`` colour (standard background 1)
    """

    def _extract_from_bgpr(bgpr_elem):
        if bgpr_elem is None:
            return None
        srgb = bgpr_elem.find(".//a:srgbClr", NS)
        if srgb is not None:
            val = srgb.get("val")
            if val:
                return str(val)
        scheme = bgpr_elem.find(".//a:schemeClr", NS)
        if scheme is not None:
            val = scheme.get("val")
            if val:
                resolved = resolve_scheme(val, theme_colors)
                if resolved:
                    return resolved
        return None

    # 1. Layout
    try:
        bg = layout.background
        result = _extract_from_bgpr(bg._element.find(".//p:bgPr", NS))
        if result:
            return result
    except Exception:
        pass

    # 2. Master
    if master is not None:
        try:
            m_bg = master.background
            result = _extract_from_bgpr(m_bg._element.find(".//p:bgPr", NS))
            if result:
                return result
        except Exception:
            pass

    # 3. Theme lt1
    return theme_colors.get("lt1")


# ── Placeholder type → human-readable label ─────────────────────────

PLACEHOLDER_LABELS: dict[int, str] = {
    1: "TITLE",
    2: "BODY",
    3: "CTR_TITLE",
    4: "SUBTITLE",
    5: "VERTICAL_TITLE",
    6: "VERTICAL_BODY",
    7: "OBJECT",
    8: "CHART",
    9: "BITMAP",
    10: "MEDIA_CLIP",
    11: "ORG_CHART",
    12: "TABLE",
    13: "SLIDE_NUMBER",
    14: "HEADER",
    15: "FOOTER",
    16: "DATE",
    17: "VERTICAL_OBJECT",
    18: "PICTURE",
    101: "SLIDE_IMAGE",
}

# Placeholder types that contain user content (ignore footer, date, slide-number)
CONTENT_PLACEHOLDER_TYPES = {
    1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 17, 18, 101,
}

# ── EMU conversion (imported from shared utils) ─────────────────────

from _ooxml_utils import emu_to_px, resolve_scheme, NS


# ── Theme extraction ─────────────────────────────────────────────────

def extract_theme_from_zip(pptx_path: Path) -> dict[str, Any]:
    """Extract theme colors and fonts from ppt/theme/*.xml inside the PPTX ZIP.

    Reads all theme files present. Colors are taken from the first theme file
    (standard). Fonts prefer a custom font scheme (name != \"Office\") over
    the built-in Office defaults, since templates often embed custom fonts in
    additional theme files.
    """
    with zipfile.ZipFile(pptx_path) as zf:
        theme_files = sorted(f for f in zf.namelist()
                             if f.startswith("ppt/theme/") and f.endswith(".xml"))
        if not theme_files:
            return {"colors": {}, "fonts": {}}
        theme_xmls = {tf: zf.read(tf) for tf in theme_files}

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    result: dict[str, Any] = {"colors": {}, "fonts": {}}

    for tf in theme_files:
        root = etree.fromstring(theme_xmls[tf])

        # Colors from first theme that has a colour scheme
        if not result["colors"]:
            clr_scheme = root.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                colors: dict[str, str] = {}
                color_names = ["dk1", "lt1", "dk2", "lt2",
                               "accent1", "accent2", "accent3",
                               "accent4", "accent5", "accent6"]
                for name in color_names:
                    elem = clr_scheme.find(f"a:{name}", ns)
                    if elem is None:
                        continue
                    srgb = elem.find("a:srgbClr", ns)
                    if srgb is not None and srgb.get("val"):
                        colors[name] = f"#{srgb.get('val')}"
                    else:
                        sys = elem.find("a:sysClr", ns)
                        if sys is not None and sys.get("lastClr"):
                            colors[name] = f"#{sys.get('lastClr')}"
                result["colors"] = colors

        # Fonts: prefer custom (non-"Office") font scheme
        font_scheme = root.find(".//a:fontScheme", ns)
        if font_scheme is not None:
            scheme_name = font_scheme.get("name", "")
            if scheme_name != "Office" or not result["fonts"]:
                major = font_scheme.find("a:majorFont", ns)
                minor = font_scheme.find("a:minorFont", ns)
                fonts: dict[str, str] = {}
                if major is not None:
                    latin = major.find("a:latin", ns)
                    if latin is not None and latin.get("typeface"):
                        fonts["major"] = latin.get("typeface")
                if minor is not None:
                    latin = minor.find("a:latin", ns)
                    if latin is not None and latin.get("typeface"):
                        fonts["minor"] = latin.get("typeface")
                if fonts:
                    result["fonts"] = fonts

    return result


# ── Custom colors extraction ─────────────────────────────────────────

def extract_custom_colors_from_zip(pptx_path: Path) -> dict[str, str]:
    """Extract custom named colours from the custClrLst element in theme XML files."""
    with zipfile.ZipFile(pptx_path) as zf:
        theme_files = sorted(f for f in zf.namelist()
                             if f.startswith("ppt/theme/") and f.endswith(".xml"))
        if not theme_files:
            return {}
        theme_xmls = {tf: zf.read(tf) for tf in theme_files}

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    custom_colors: dict[str, str] = {}

    for tf in theme_files:
        root = etree.fromstring(theme_xmls[tf])
        cust_clr_lst = root.find(".//a:custClrLst", ns)
        if cust_clr_lst is not None:
            for cust_clr in cust_clr_lst.findall("a:custClr", ns):
                name = cust_clr.get("name")
                srgb = cust_clr.find("a:srgbClr", ns)
                if name and srgb is not None and srgb.get("val"):
                    custom_colors[name] = f"#{srgb.get('val')}"

    return custom_colors


# ── Image extraction from shape ──────────────────────────────────────

# Formats that browsers natively support in <img> tags
_WEB_COMPATIBLE_EXTS = frozenset({"png", "jpg", "jpeg", "gif", "svg", "webp", "ico"})

# Normalise MIME-type-derived extensions (e.g. "svg+xml" → "svg")
_EXT_TO_SAFE = {
    "svg+xml": "svg",
}


def _convert_to_png(blob: bytes, source_ext: str) -> bytes | None:
    """Attempt to convert an image blob to PNG using Pillow.

    Returns the PNG bytes on success, or None if conversion is not possible
    (e.g. WMF/EMF without platform support).
    """
    import io
    from PIL import Image, UnidentifiedImageError

    try:
        img = Image.open(io.BytesIO(blob))
        if img.format in ("EMF", "WMF", "EMF+"):
            # These formats require platform-specific rasterisers (Windows GDI).
            # On macOS/Linux they open but cannot be fully loaded.
            return None
        img = img.convert("RGBA")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return None


def _gen_web_compatible(
    blob: bytes, ext: str, assets_dir: Path, basename: str,
    seen_hashes: dict[str, str],
) -> tuple[str, bool]:
    """Generate a web-compatible image file from an extracted blob.

    If the extension is already web-compatible, saves the blob as-is.
    Otherwise, attempts conversion to PNG.  Deduplicates via SHA-256.

    Returns ``(filename, web_compatible)`` where *filename* is the file
    name relative to the assets directory and *web_compatible* indicates
    whether the saved file can be rendered in a browser ``<img>``.
    """
    h = hashlib.sha256(blob).hexdigest()[:8]

    if h in seen_hashes:
        return seen_hashes[h], seen_hashes.get(h + ":wc", True)

    assets_dir.mkdir(parents=True, exist_ok=True)

    ext_lower = ext.lower() if ext else "png"

    if ext_lower in _WEB_COMPATIBLE_EXTS:
        filename = f"{basename}_{h}.{ext_lower}"
        (assets_dir / filename).write_bytes(blob)
        seen_hashes[h] = filename
        seen_hashes[h + ":wc"] = True
        return filename, True

    # Try conversion to PNG
    png = _convert_to_png(blob, ext_lower)
    if png:
        filename = f"{basename}_{h}.png"
        (assets_dir / filename).write_bytes(png)
        seen_hashes[h] = filename
        seen_hashes[h + ":wc"] = True
        return filename, True

    # Save original as fallback
    filename = f"{basename}_{h}.{ext_lower}"
    (assets_dir / filename).write_bytes(blob)
    seen_hashes[h] = filename
    seen_hashes[h + ":wc"] = False
    return filename, False


def _extract_image_from_shape(shape, layout=None, verbose: bool = False) -> tuple[bytes | None, str | None, str | None]:
    """Extract the embedded image from a PICTURE shape, using python-pptx's
    internal relationship chain. Works on any shape in the slide master or
    any layout — no hardcoded names or paths.

    Handles both regular embedded images (r:embed on a:blip) and SVG
    extensions (asvg:svgBlip inside a:extLst).

    When ``layout`` is provided (a pptx SlideLayout), it is used as fallback
    for resolving relationships when ``shape.part.related_parts`` is empty.

    If ``verbose``, prints diagnostic trace to stderr on failure.

    Returns ``(blob, extension, content_type)`` or ``(None, None, None)``.
    """
    import sys as _sys

    # High-level API (works for shapes in slide layouts)
    try:
        blob = shape.image.blob
        if verbose:
            print(f"  [img] high-level API OK: {len(blob)}b", file=_sys.stderr)
        return blob, shape.image.ext or "png", shape.image.content_type or "image/png"
    except Exception as exc:
        if verbose:
            print(f"  [img] high-level API failed: {exc}", file=_sys.stderr)

    _NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    _NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    _NS_ASVG = "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}"

    _MIME_BY_EXT = {
        "svg": "image/svg+xml", "svg+xml": "image/svg+xml", "png": "image/png",
        "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "gif": "image/gif", "bmp": "image/bmp",
        "tiff": "image/tiff", "tif": "image/tiff",
        "emf": "image/emf", "wmf": "image/wmf",
    }

    def _resolve(blip_elem) -> tuple[bytes | None, str | None, str | None] | None:
        """Try to resolve an image by walking blip → r:embed → related_parts."""
        r_embed = blip_elem.get(_NS_R + "embed")
        if not r_embed:
            if verbose:
                print(f"  [img]   no r:embed on blip elem", file=_sys.stderr)
            return None
        if verbose:
            print(f"  [img]   r:embed={r_embed}", file=_sys.stderr)
        # Strategy 1: shape.part.related_parts
        rels = getattr(shape.part, "related_parts", {})
        image_part = rels.get(r_embed)
        if image_part is None:
            if verbose:
                print(f"  [img]   shape.part.related_parts missing (keys={list(rels.keys())[:5]})", file=_sys.stderr)
        # Strategy 2: layout.part.related_part (handles shapes on slide layouts)
        if image_part is None and layout is not None:
            try:
                image_part = layout.part.related_part(r_embed)
                if verbose:
                    print(f"  [img]   layout.part.related_part OK: {type(image_part).__name__}", file=_sys.stderr)
            except Exception as exc:
                if verbose:
                    print(f"  [img]   layout.part.related_part failed: {exc}", file=_sys.stderr)
        if image_part is None:
            if verbose:
                print(f"  [img]   resolution FAILED", file=_sys.stderr)
            return None
        blob = image_part.blob
        content_type = getattr(image_part, "content_type", "") or "image/png"
        ext = content_type.rsplit("/", 1)[-1] if "/" in content_type else "png"
        ext = _EXT_TO_SAFE.get(ext, ext)
        ct = _MIME_BY_EXT.get(ext, content_type)
        if verbose:
            print(f"  [img]   resolved: {len(blob)}b ct={ct}", file=_sys.stderr)
        return blob, ext, ct

    try:
        blip = shape._element.find(f".//{_NS_A}blip")
        if blip is None:
            if verbose:
                print(f"  [img] a:blip not found in shape._element", file=_sys.stderr)
            return None, None, None

        if verbose:
            blip_direct = blip.get(_NS_R + "embed")
            svg_sub = blip.find(f".//{_NS_ASVG}svgBlip") is not None
            print(f"  [img] blip found, direct r:embed={blip_direct or 'MISSING'}, has_svg={svg_sub}", file=_sys.stderr)

        # 1) Direct a:blip/@r:embed
        result = _resolve(blip)
        if result is not None:
            return result

        # 2) SVG blip inside a:extLst/a:ext/asvg:svgBlip
        svg_blip = blip.find(f".//{_NS_ASVG}svgBlip")
        if svg_blip is not None:
            if verbose:
                print(f"  [img] trying SVG blip resolution", file=_sys.stderr)
            result = _resolve(svg_blip)
            if result is not None:
                return result
        elif verbose:
            print(f"  [img] no SVG blip found under a:blip", file=_sys.stderr)

        if verbose:
            print(f"  [img] ALL strategies failed", file=_sys.stderr)
        return None, None, None
    except Exception as exc:
        if verbose:
            print(f"  [img] FATAL: {exc}", file=_sys.stderr)
        return None, None, None


# ── Slide master element extraction ──────────────────────────────────

def extract_master_elements(
    slide_master,
    assets_dir: Path,
    theme_colors: dict[str, str] | None = None,
) -> dict[str, Any]:
    """Extract reserved areas (footer, slide-number, logo) from the slide master.

    Logo images are written to ``assets_dir/`` and embedded as base64 data URIs
    so the HTML skeleton is self-contained.

    When ``theme_colors`` is provided, text style information (font size, colour,
    alignment, etc.) is extracted from each shape's ``p:txBody`` XML element.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    theme_colors = theme_colors or {}
    reserved_areas: list[dict[str, Any]] = []

    for shape in slide_master.shapes:
        if shape.is_placeholder:
            ph_type = int(shape.placeholder_format.type)
            if ph_type in (15, 13):  # FOOTER, SLIDE_NUMBER
                text = shape.text_frame.text if shape.has_text_frame else ""
                label = PLACEHOLDER_LABELS.get(ph_type, f"TYPE_{ph_type}")
                entry: dict[str, Any] = {
                    "type": label,
                    "x": emu_to_px(shape.left),
                    "y": emu_to_px(shape.top),
                    "w": emu_to_px(shape.width),
                    "h": emu_to_px(shape.height),
                    "text": text,
                }
                if shape.has_text_frame and shape._element is not None:
                    ts = extract_text_style(shape._element, theme_colors)
                    if ts:
                        entry["text_style"] = ts
                reserved_areas.append(entry)
        else:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_blob, ext, content_type = _extract_image_from_shape(shape, layout=slide_master)
                if image_blob is None:
                    print(f"  Warning: Could not extract master picture '{shape.name}': no image data")
                    continue

                # SHA-256 dedup (same image on multiple masters → one file)
                h = hashlib.sha256(image_blob).hexdigest()[:8]
                ext_safe = _EXT_TO_SAFE.get((ext or "png").lower(), (ext or "png").lower())
                logo_filename = f"logo_{h}.{ext_safe}"
                assets_dir.mkdir(parents=True, exist_ok=True)
                if not (assets_dir / logo_filename).exists():
                    (assets_dir / logo_filename).write_bytes(image_blob)

                b64 = base64.b64encode(image_blob).decode("ascii")
                logo_src = f"data:{content_type or 'image/' + ext_safe};base64,{b64}"

                reserved_areas.append({
                    "type": "LOGO",
                    "x": emu_to_px(shape.left),
                    "y": emu_to_px(shape.top),
                    "w": emu_to_px(shape.width),
                    "h": emu_to_px(shape.height),
                    "text": "",
                    "logo_src": logo_src,
                })
                print(f"  Extracted master logo: {logo_filename} ({len(image_blob)} bytes)")

    return {"reserved_areas": reserved_areas}


# ── Layout shape extraction ──────────────────────────────────────────

def extract_layout_shapes(
    layout,
    theme_colors: dict[str, str] | None = None,
    assets_dir: Path | None = None,
    seen_image_hashes: dict[str, Any] | None = None,
    verbose: bool = False,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Extract content placeholders (with default text + text style) and
    immutable non-PH shapes.

    When ``assets_dir`` is provided, PICTURE shapes have their image
    extracted to ``assets_dir/`` with SHA-256 deduplication.
    ``seen_image_hashes`` is a shared dict (hash → filename) to avoid
    saving duplicate image files across multiple layout calls.

    ``layout`` is also passed to ``_extract_image_from_shape`` for
    relationship resolution of SVG images stored on slide layouts.

    Returns ``(placeholders, immutable_shapes)``.

    When ``theme_colors`` is provided, full text formatting is extracted from
    each shape's underlying ``p:txBody`` XML element.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    theme_colors = theme_colors or {}
    seen = seen_image_hashes or {}

    placeholders: list[dict[str, Any]] = []
    immutable_shapes: list[dict[str, Any]] = []

    for shape in layout.shapes:
        if shape.is_placeholder:
            ptype = int(shape.placeholder_format.type)
            if ptype not in CONTENT_PLACEHOLDER_TYPES:
                continue
            label = PLACEHOLDER_LABELS.get(ptype, f"TYPE_{ptype}")

            default_text = ""
            if shape.has_text_frame:
                default_text = shape.text_frame.text.strip()

            ph_entry: dict[str, Any] = {
                "type": label,
                "x": emu_to_px(shape.left),
                "y": emu_to_px(shape.top),
                "w": emu_to_px(shape.width),
                "h": emu_to_px(shape.height),
                "idx": int(shape.placeholder_format.idx),
                "name": shape.name,
            }
            if default_text:
                ph_entry["default_text"] = default_text
            if shape.has_text_frame and shape._element is not None:
                ts = extract_text_style(shape._element, theme_colors)
                if ts:
                    ph_entry["text_style"] = ts
            placeholders.append(ph_entry)
        else:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                children = flatten_group_shape(shape, theme_colors)
                immutable_shapes.extend(children)
            else:
                sh_entry = _make_shape_entry(
                    shape, shape.shape_type,
                    emu_to_px(shape.left), emu_to_px(shape.top),
                    emu_to_px(shape.width), emu_to_px(shape.height),
                    theme_colors,
                )
                if sh_entry:
                    # Extract image from PICTURE shapes with dedup
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and assets_dir is not None:
                        img_blob, ext, _ct = _extract_image_from_shape(shape, layout=layout, verbose=verbose)
                        if img_blob:
                            filename, web_ok = _gen_web_compatible(
                                img_blob, ext or "png", assets_dir, "logo", seen
                            )
                            sh_entry["image_src"] = f"../assets/{filename}"
                            if not web_ok:
                                sh_entry["web_compatible"] = False

                    # Also grab text style for shapes with text
                    if shape.has_text_frame and shape._element is not None:
                        text = shape.text_frame.text.strip()
                        if text and "text" not in sh_entry:
                            sh_entry["text"] = text
                        ts = extract_text_style(shape._element, theme_colors)
                        if ts:
                            sh_entry["text_style"] = ts
                    immutable_shapes.append(sh_entry)

    return placeholders, immutable_shapes

def generate_colors_md(theme: dict[str, Any], custom_colors: dict[str, str] | None = None) -> str:
    """Generate colors.md in a format compatible with slide_compile_theme."""
    colors = theme.get("colors", {})
    custom_colors = custom_colors or {}

    token_table = [
        "| Token Name | Hex Code |",
        "| --- | --- |",
    ]
    for name in ["dk1", "lt1", "dk2", "lt2",
                 "accent1", "accent2", "accent3",
                 "accent4", "accent5", "accent6"]:
        if name in colors:
            token_table.append(f"| `{name}` | `{colors[name]}` |")

    # Map OOXML colors to semantic roles
    accent = colors.get("accent1", "#4472C4")
    accent_secondary = colors.get("accent2", accent)
    bg_light = colors.get("lt1", "#FFFFFF")
    text_dark = colors.get("dk1", "#000000")
    bg_dark = colors.get("dk2", "#44546A")

    lines = [
        "---",
        "trigger: always_on",
        "---",
        "",
        "# Color System",
        "",
        "## 1. Global Color Tokens",
        "",
        *token_table,
        "",
        "## 2. Custom Named Colors",
        "",
    ]

    if custom_colors:
        lines += [
            "| Name | Hex Code |",
            "| --- | --- |",
        ]
        for name, hex_code in sorted(custom_colors.items()):
            lines.append(f"| `{name}` | `{hex_code}` |")
        lines.append("")

    lines += [
        "## 3. Theme Architecture",
        "",
        f"- Background Light: `{bg_light}`",
        f"- Background Dark: `{bg_dark}`",
        f"- Text Primary: `{text_dark}`",
        f"- Accent: `{accent}`",
        f"- Accent Secondary: `{accent_secondary}`",
        f"- Border: `#D9D9D9`",
        "",
        "## 4. CSS Implementation",
        "",
        "```css",
        ":root {",
        f"  --c-accent: {accent};",
        f"  --c-accent-secondary: {accent_secondary};",
        f"  --c-bg: {bg_light};",
        f"  --c-text: {text_dark};",
        "}",
        "```",
    ]
    return "\n".join(lines) + "\n"


def generate_typography_md(theme: dict[str, Any]) -> str:
    """Generate typography.md in a format compatible with slide_compile_theme."""
    fonts = theme.get("fonts", {})
    heading = fonts.get("major", "Calibri")
    body = fonts.get("minor", heading)

    return "\n".join([
        "---",
        "trigger: always_on",
        "---",
        "",
        "# Typography",
        "",
        "## 1. Font Families & Usage",
        "",
        "| Family | Variant | Usage | Style Rules |",
        "| --- | --- | --- | --- |",
        f"| **{heading}** | **Regular** | Headlines: titles and subtitles. | Use for headings. |",
        f"| **{body}** | **Regular** | Body copy: long text, paragraphs. | Use for body text. |",
        "",
        "## 2. Font Files",
        "",
        "No local font files — Google Fonts CDN fallback will be used.",
        "",
    ]) + "\n"





# ── Fingerprint ──────────────────────────────────────────────────────

def compute_fingerprint(pptx_path: Path, layout_count: int) -> dict[str, Any]:
    """Compute SHA-256 fingerprint of the template file."""
    sha = hashlib.sha256()
    with open(pptx_path, "rb") as f:
        while True:
            chunk = f.read(65536)
            if not chunk:
                break
            sha.update(chunk)
    return {
        "sha256": sha.hexdigest(),
        "layout_count": layout_count,
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def read_fingerprint(output_dir: Path) -> dict[str, Any] | None:
    """Read existing fingerprint, or None."""
    fp_path = output_dir / ".template-fingerprint"
    if not fp_path.exists():
        return None
    try:
        return json.loads(fp_path.read_text("utf-8"))
    except Exception:
        return None


def write_fingerprint(output_dir: Path, fp: dict[str, Any]) -> None:
    """Write fingerprint file."""
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / ".template-fingerprint").write_text(
        json.dumps(fp, indent=2) + "\n", "utf-8"
    )


# ── Main analyzer ────────────────────────────────────────────────────

def analyze_template(
    pptx_path: Path,
    design_name: str,
    output_dir: Path,
    *,
    force: bool = False,
    verbose: bool = False,
) -> str:
    """Analyze a corporate .pptx template and produce a design system."""

    if not pptx_path.exists():
        return f"Error: Template file not found: {pptx_path}"
    if pptx_path.suffix.lower() != ".pptx":
        return f"Error: Must be a .pptx file, got: {pptx_path.suffix}"

    design_dir = output_dir / design_name

    # ── Check fingerprint for idempotency ──────────────────────
    if not force:
        existing = read_fingerprint(design_dir)
        if existing:
            current = compute_fingerprint(pptx_path, 0)  # count not needed for check
            if current["sha256"] == existing["sha256"]:
                return (
                    f"Template '{design_name}' already analyzed "
                    f"({existing['layout_count']} layouts). Skipping.\n"
                    f"Use --force to re-analyze."
                )
            else:
                print(
                    f"Template changed (sha256 mismatch). Re-analyzing "
                    f"(was {existing.get('layout_count', '?')} layouts)..."
                )

    # ── Load presentation ──────────────────────────────────────
    prs = Presentation(str(pptx_path))
    sm = prs.slide_masters[0]
    layouts = list(sm.slide_layouts)
    total = len(layouts)
    print(f"Analyzing {total} slide layouts from template...")

    # ── Extract theme ──────────────────────────────────────────
    theme = extract_theme_from_zip(pptx_path)
    custom_colors = extract_custom_colors_from_zip(pptx_path)
    theme_colors = theme.get("colors", {})
    theme_fonts = dict(theme.get("fonts", {}))

    # ── Extract master p:txStyles from ZIP ────────────────────
    master_text_styles: dict[str, Any] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        master_xml_files = sorted(f for f in zf.namelist()
                                   if f.startswith("ppt/slideMasters/") and f.endswith(".xml"))
        if master_xml_files:
            master_xml_bytes = zf.read(master_xml_files[0])
            master_text_styles = extract_master_text_styles(master_xml_bytes, theme_colors)

    # ── Create output directories ──────────────────────────────
    design_dir.mkdir(parents=True, exist_ok=True)
    layouts_dir = design_dir / "layouts"
    layouts_dir.mkdir(parents=True, exist_ok=True)
    assets_dir = design_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    screenshots_dir = design_dir / "screenshots"
    screenshots_dir.mkdir(parents=True, exist_ok=True)

    # ── Extract slide master elements ─────────────────────────
    master = extract_master_elements(sm, assets_dir, theme_colors)
    reserved_areas = master["reserved_areas"]
    print(f"  Master: {len(reserved_areas)} reserved areas (footer/logo/slide-number)")

    # ── Enumerate layouts ──────────────────────────────────────
    layout_entries: dict[str, dict[str, Any]] = {}
    layout_categories: dict[str, list[str]] = {}
    seen_image_hashes: dict[str, Any] = {}

    for idx, layout in enumerate(layouts):
        name = layout.name or f"Layout_{idx + 1}"
        safe_key = name.lower().replace(" ", "_").replace("/", "_")[:60]

        ph_list, immutable_shapes = extract_layout_shapes(layout, theme_colors, assets_dir, seen_image_hashes, verbose=verbose)

        # ── Merge master text styles into per-shape entries ──────
        for ph in ph_list:
            ts = ph.get("text_style")
            if ts:
                ph["text_style"] = _merge_master_text_style(
                    ts, ph.get("type", ""), master_text_styles, theme_fonts
                )
        for sh in immutable_shapes:
            ts = sh.get("text_style")
            if ts:
                sh["text_style"] = _merge_master_text_style(
                    ts, sh.get("shape_type", ""), master_text_styles, theme_fonts
                )

        category = classify_layout(name, ph_list)

        bg_color = _detect_layout_bg(layout, master, theme_colors)

        layout_entries[safe_key] = {
            "file": f"{safe_key}.html",
            "category": category,
            "original_name": name,
            "placeholders": ph_list,
            "bg_color": bg_color,
            "immutable_shapes": immutable_shapes,
        }

        layout_categories.setdefault(category, []).append(safe_key)

    # ── Write colors.md and typography.md ──────────────────────
    (design_dir / "colors.md").write_text(generate_colors_md(theme, custom_colors), "utf-8")
    (design_dir / "typography.md").write_text(generate_typography_md(theme), "utf-8")

    # ── Write layout skeletons ─────────────────────────────────
    for key, entry in layout_entries.items():
        html = generate_skeleton_html(
            layout_name=entry["original_name"],
            placeholders=entry["placeholders"],
            bg_color=entry.get("bg_color"),
            reserved_areas=reserved_areas,
            immutable_shapes=entry.get("immutable_shapes"),
            theme_fonts=theme_fonts,
        )
        (layouts_dir / entry["file"]).write_text(html, "utf-8")

    # ── Write index.json ───────────────────────────────────────
    index = {
        "design_system": design_name,
        "total_layouts": total,
        "theme_colors": theme_colors,
        "theme_fonts": theme_fonts,
        "master_elements": {
            "reserved_areas": reserved_areas,
            "text_styles": master_text_styles,
        },
        "custom_colors": custom_colors,
        "categories": {
            cat: {
                "count": len(names),
                "layouts": names,
            }
            for cat, names in sorted(layout_categories.items())
        },
        "layouts": layout_entries,
        "bg_colors": {
            safe_key: entry.get("bg_color")
            for safe_key, entry in layout_entries.items()
            if entry.get("bg_color")
        },
    }
    (layouts_dir / "index.json").write_text(
        json.dumps(index, indent=2, ensure_ascii=False) + "\n", "utf-8"
    )

    # ── Write fingerprint ──────────────────────────────────────
    fp = compute_fingerprint(pptx_path, total)
    write_fingerprint(design_dir, fp)

    # ── Summary ────────────────────────────────────────────────
    immutable_count = sum(
        1 for e in layout_entries.values() if e.get("immutable_shapes")
    )
    lines = [
        f"Template analyzed: {design_name}",
        f"  Layouts: {total}",
        f"  Categories: {len(layout_categories)}",
        f"  Design system: {design_dir}",
        f"  Theme colors: {len(theme.get('colors', {}))} extracted + {len(custom_colors)} custom",
        f"  Theme fonts: {theme.get('fonts', {})}",
        f"  Master reserved: {len(reserved_areas)} items (footer/logo/slide-number) + {len(master_text_styles)} text style sections",
        f"  Immutable shapes: {immutable_count} layouts have non-PH decorative shapes",
        f"  Fingerprint: {fp['sha256'][:12]}...",
    ]
    return "\n".join(lines)


# ── CLI ──────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Analyze a corporate .pptx template and produce a design system."
    )
    parser.add_argument("--input", required=True, help="Path to template .pptx file")
    parser.add_argument("--design-name", required=True, help="Name for the design system directory")
    parser.add_argument("--output-dir", required=True, help="Path to .opencode/office/slides/design/ directory")
    parser.add_argument("--force", action="store_true", help="Re-analyze even if fingerprint matches")
    parser.add_argument("--verbose", action="store_true", help="Print diagnostic output for image extraction failures")
    args = parser.parse_args()

    result = analyze_template(
        pptx_path=Path(args.input),
        design_name=args.design_name,
        output_dir=Path(args.output_dir),
        force=args.force,
        verbose=args.verbose,
    )
    print(result)


if __name__ == "__main__":
    main()
