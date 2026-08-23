"""
Assemble a .pptx presentation from a corporate template, populating
placeholder content from authored HTML slide files.

Usage:
  python assemble_from_template.py \\
    --template corporate.pptx \\
    --slides slide_01.html slide_02.html ... \\
    --layout-map '{"slide_01.html": "title_centered_01", ...}' \\
    --output output.pptx
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt, Emu  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore


# ── Placeholder type label → PP_PLACEHOLDER enum value ───────────────

LABEL_TO_PPTX_TYPE: dict[str, Any] = {
    "TITLE": 1,
    "BODY": 2,
    "CTR_TITLE": 3,
    "SUBTITLE": 4,
    "VERTICAL_TITLE": 5,
    "VERTICAL_BODY": 6,
    "OBJECT": 7,
    "CHART": 8,
    "BITMAP": 9,
    "MEDIA_CLIP": 10,
    "ORG_CHART": 11,
    "TABLE": 12,
    "PICTURE": 18,
    "SLIDE_IMAGE": 101,
}


# ── HTML → python-pptx text extraction ────────────────────────

def hex_to_rgb(hex_color: str) -> RGBColor | None:
    """Convert '#rrggbb' or '#rgb' to RGBColor."""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    if len(hex_color) != 6:
        return None
    try:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
    except ValueError:
        return None


def _parse_inline_style(style_str: str) -> dict[str, str]:
    """Parse inline style='...' into a dict of property: value."""
    import re
    matches = re.findall(r'([\w-]+)\s*:\s*([^;]+)', style_str)
    return {k.strip(): v.strip() for k, v in matches}


def _apply_style_to_run(run, css_props: dict[str, str]) -> None:
    """Apply CSS properties from an inline style dict to a pptx run."""
    import re as _re

    # font-size: value is e.g. "14px"
    fs_val = css_props.get("font-size", "")
    fs_match = _re.match(r'(\d+)', fs_val)
    if fs_match:
        run.font.size = Pt(int(fs_match.group(1)))

    # font-family: value is e.g. "'Baikal', sans-serif"
    ff_val = css_props.get("font-family", "")
    ff_match = _re.search(r'["\']?([A-Za-z0-9 _-]+)', ff_val)
    if ff_match:
        run.font.name = ff_match.group(1).strip()

    # font-weight: "bold" or "700" etc.
    fw_val = css_props.get("font-weight", "").strip().lower()
    if fw_val in ("bold", "bolder") or (fw_val.isdigit() and int(fw_val) >= 700):
        run.font.bold = True

    # font-style: italic
    if "italic" in css_props.get("font-style", "").lower():
        run.font.italic = True

    # text-decoration: underline
    if "underline" in css_props.get("text-decoration", ""):
        run.font.underline = True

    # color: value is e.g. "#e60000"
    color_val = css_props.get("color", "")
    color_match = _re.search(r'(#[0-9a-fA-F]{3,6})', color_val)
    if color_match:
        rgb = hex_to_rgb(color_match.group(1))
        if rgb:
            run.font.color.rgb = rgb


def _apply_style_to_paragraph(paragraph, css_props: dict[str, str]) -> None:
    """Apply CSS properties to a paragraph."""
    align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                 "left": PP_ALIGN.LEFT, "justify": PP_ALIGN.JUSTIFY}
    align_val = css_props.get("text-align", "").strip().lower()
    if align_val in align_map:
        paragraph.alignment = align_map[align_val]


def _extract_text_from_element(
    el,
    text_frame,
    level: int = 0,
) -> None:
    """Recursively extract text content from an HTML element into pptx paragraphs."""
    # Text-bearing tags (paragraph-level)
    text_tags = {"p", "h1", "h2", "h3", "h4", "h5", "h6", "li"}
    # Container tags to recurse into
    container_tags = {"div", "section", "article", "main", "header", "footer", "nav"}
    # Inline / list tags
    inline_tags = {"b", "i", "strong", "em", "span", "a", "br"}

    if not hasattr(el, "name") or el.name is None:
        text = el.string
        if text and text.strip():
            run = text_frame.paragraphs[-1].add_run()
            run.text = text
        return

    tag = el.name.lower() if el.name else ""

    # Parse this element's inline style once
    style_str = el.get("style", "") if hasattr(el, "get") else ""
    css_props = _parse_inline_style(style_str)

    if tag in text_tags:
        p = text_frame.add_paragraph()
        p.text = ""
        _apply_style_to_paragraph(p, css_props)
        if tag.startswith("h"):
            try:
                lvl = int(tag[1])
                sizes = {1: Pt(32), 2: Pt(24), 3: Pt(20), 4: Pt(18), 5: Pt(16), 6: Pt(14)}
                p.font.size = sizes.get(lvl, Pt(16))
                p.font.bold = True
            except (ValueError, TypeError):
                pass
        if tag == "li":
            p.level = level

    for child in el.children:
        if hasattr(child, "name") and child.name == "br":
            text_frame.add_paragraph().text = ""
            continue

        if hasattr(child, "name") and child.name in ("ul", "ol"):
            _extract_text_from_element(child, text_frame, level + 1)
            continue

        if hasattr(child, "name") and child.name and (child.name.lower() in text_tags
                                       or child.name.lower() in inline_tags
                                       or child.name.lower() in container_tags
                                       or child.name.lower() in ("ul", "ol")):
            _extract_text_from_element(child, text_frame, level)
            continue

        # Bare text node
        if hasattr(child, "string") and child.string:
            txt = child.string.strip()
            if txt:
                if len(text_frame.paragraphs) == 0:
                    text_frame.add_paragraph().text = ""
                run = text_frame.paragraphs[-1].add_run()
                run.text = txt

                # Apply inline styles
                parent = child.parent
                if parent and hasattr(parent, "name"):
                    pname = parent.name.lower() if parent.name else ""
                    if pname in ("b", "strong"):
                        run.font.bold = True
                    if pname in ("i", "em"):
                        run.font.italic = True
                    if parent.get("style"):
                        parent_css = _parse_inline_style(parent.get("style", ""))
                        _apply_style_to_run(run, parent_css)

                # Also apply this element's own CSS (for direct text nodes in styled divs)
                if css_props and not (tag in text_tags or tag in container_tags):
                    _apply_style_to_run(run, css_props)


def populate_placeholder(placeholder, html_element) -> None:
    """Populate a PPTX placeholder with content from an HTML element."""
    tf = placeholder.text_frame
    tf.clear()

    _extract_text_from_element(html_element, tf)

    # If nothing was extracted, add empty paragraph to avoid corrupt text frame
    if len(tf.paragraphs) == 0:
        tf.add_paragraph().text = ""


# ── Image extraction from HTML ───────────────────────────────────────

def extract_images_from_placeholder(html_element, slide_dir: Path) -> list[Path]:
    """Find <img> tags within a placeholder element, resolve paths."""
    images: list[Path] = []
    for img in html_element.find_all("img"):
        src = img.get("src", "")
        if not src:
            continue
        img_path = slide_dir / src.lstrip("./")
        if img_path.exists():
            images.append(img_path)
    return images


# ── Layout resolution ────────────────────────────────────────────────

def find_layout_by_name(prs: Presentation, name: str):
    """Find a slide layout by its name in the presentation."""
    # Exact match
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout

    # Normalize both sides: lowercased, underscored, special chars replaced
    def _normalize(s: str) -> str:
        return s.lower().replace("_&_", "_and_").replace(" & ", "_and_").replace("& ", "_and_").replace(" &", "_and_") \
                .replace("&", "_and_").replace("+", "_plus_").replace("-", "_") \
                .replace(" ", "_").replace("/", "_")

    safe = _normalize(name)
    for layout in prs.slide_layouts:
        if _normalize(layout.name) == safe:
            return layout
    return None


# ── Main assembler ───────────────────────────────────────────────────

def assemble_presentation(
    template_path: Path,
    slide_paths: list[Path],
    layout_map: dict[str, str],
    output_path: Path,
) -> str:
    """Assemble a PPTX from a template, populating placeholders from HTML."""

    if not template_path.exists():
        return f"Error: Template not found: {template_path}"

    # Copy template to preserve masters
    if template_path != output_path:
        shutil.copy2(template_path, output_path)
        prs = Presentation(str(output_path))
    else:
        prs = Presentation(str(template_path))

    # Delete all existing slides from the template copy
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    errors: list[str] = []

    for slide_path in slide_paths:
        slide_name = slide_path.name
        layout_name = layout_map.get(slide_name)
        if not layout_name:
            errors.append(f"Warning: {slide_name} not in layout map, skipping")
            continue

        layout = find_layout_by_name(prs, layout_name)
        if layout is None:
            errors.append(
                f"Error: Layout '{layout_name}' not found in template "
                f"(slide: {slide_name})"
            )
            continue

        # Create slide from the layout
        slide = prs.slides.add_slide(layout)

        # Parse authored HTML
        if not slide_path.exists():
            errors.append(f"Warning: Slide file not found: {slide_path}")
            continue

        html_text = slide_path.read_text("utf-8")
        soup = BeautifulSoup(html_text, "html.parser")

        # Find all [data-placeholder] elements
        ph_elements = soup.select("[data-placeholder]")
        if not ph_elements:
            errors.append(f"Warning: No data-placeholder elements in {slide_name}")

        slide_dir = slide_path.parent

        for ph_el in ph_elements:
            ph_type = ph_el.get("data-placeholder", "")

            if ph_type not in LABEL_TO_PPTX_TYPE:
                errors.append(f"Warning: Unknown placeholder type '{ph_type}' in {slide_name}")
                continue

            pptx_type = LABEL_TO_PPTX_TYPE[ph_type]

            # Parse HTML placeholder position (in pixels, from inline style)
            ph_style = ph_el.get("style", "")
            import re as _re
            pos_match = _re.findall(r'(left|top|width|height)\s*:\s*(\d+)px', ph_style)
            ph_pos = {k: int(v) for k, v in pos_match}

            # ── Primary match: data-idx (added by generate_contracts.py) ──
            html_idx = ph_el.get("data-idx", "")
            target_ph = None
            PX_TO_EMU = 9525  # 1px ≈ 9525 EMU (at 96 DPI)
            TOLERANCE = 20000  # ~2px tolerance

            if html_idx:
                try:
                    html_idx_int = int(html_idx)
                    for ph in slide.placeholders:
                        ptype_int = int(ph.placeholder_format.type)
                        if ptype_int == pptx_type and int(ph.placeholder_format.idx) == html_idx_int:
                            target_ph = ph
                            break
                except ValueError:
                    pass

            # ── Fallback: position-based proximity match ──
            if target_ph is None:
                best_dist = float('inf')
                for ph in slide.placeholders:
                    if int(ph.placeholder_format.type) != pptx_type:
                        continue
                    if ph_pos:
                        ph_px_left = ph_pos.get('left', 0)
                        ph_px_top = ph_pos.get('top', 0)
                        ph_emu_left = ph.left or 0
                        ph_emu_top = ph.top or 0
                        dist = abs(ph_emu_left - ph_px_left * PX_TO_EMU) + abs(ph_emu_top - ph_px_top * PX_TO_EMU)
                        if dist < best_dist:
                            best_dist = dist
                            target_ph = ph
                    else:
                        target_ph = ph
                        break

                # If proximity match too far, reset for type-only fallback
                if ph_pos and best_dist > TOLERANCE * 2:
                    target_ph = None

            # ── Final fallback: first matching type ──
            if target_ph is None:
                for ph in slide.placeholders:
                    if int(ph.placeholder_format.type) == pptx_type:
                        target_ph = ph
                        break

            if target_ph is None:
                errors.append(
                    f"Warning: No placeholder of type '{ph_type}' on slide "
                    f"from layout '{layout_name}' ({slide_name})"
                )
                continue

            # Handle images in PICTURE/OBJECT placeholders
            if ph_type in ("PICTURE", "OBJECT", "BITMAP", "SLIDE_IMAGE"):
                images = extract_images_from_placeholder(ph_el, slide_dir)
                if images:
                    try:
                        target_ph.insert_picture(str(images[0]))
                    except Exception as exc:
                        errors.append(f"Error inserting image in {slide_name}: {exc}")
                else:
                    populate_placeholder(target_ph, ph_el)
            else:
                populate_placeholder(target_ph, ph_el)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    summary = [
        f"Presentation assembled: {output_path}",
        f"Slides: {len(slide_paths)}",
    ]
    if errors:
        summary.append("Issues:")
        summary.extend(f"  {e}" for e in errors)

    return "\n".join(summary)


# ── CLI ──────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Assemble PPTX from corporate template using HTML slide content."
    )
    parser.add_argument("--template", required=True, help="Path to corporate .pptx template")
    parser.add_argument("--slides", nargs="+", required=True, help="HTML slide files in order")
    parser.add_argument("--layout-map", required=True, help="JSON mapping slide filenames to layout names")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    layout_map: dict[str, str] = json.loads(args.layout_map)
    slide_paths = [Path(s) for s in args.slides]

    result = assemble_presentation(
        template_path=Path(args.template),
        slide_paths=slide_paths,
        layout_map=layout_map,
        output_path=Path(args.output),
    )
    print(result)


if __name__ == "__main__":
    main()
